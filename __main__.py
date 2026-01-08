# Serverless function for PowerPoint generation
import json
import ast
import base64
import io
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import os
import traceback

# Power Automate endpoint URL
POWER_AUTOMATE_ENDPOINT = "https://d2f5de0d12414bd1a0bbd898f8452b.88.environment.api.powerplatform.com:443/powerautomate/automations/direct/workflows/400ab3f3afef4518bbaf6d4f7a8a6c8f/triggers/manual/paths/invoke?api-version=1&sp=%2Ftriggers%2Fmanual%2Frun&sv=1.0&sig=AQbZ1utjacB3WOMaHAsdra8hRfVazdsdxsD67wrTx3s"

# ============================================================================
# DATA FIELD DEFINITIONS
# ============================================================================

SLIDE_1_KEYS = [
    'Project Name',
    'PMD Name',
    'Overall Status',
    'Staffing Status',
    'Scope Status',
    'Project Governance Status',
    'Escalation Management Status1',
    'Reason for Amber / Red',
    'Total FTEs',
    'Project Highlights',
    'Up-sell / Cross-sell Opportunities',
    'Up-sell / Cross-sell Details'
]

SLIDE_2_KEYS = [
    'Project Name',
    'Tool',
    'Description',
    'Steps to Reproduce',
    'Module',
    'Priority',
    'Owner',
    'SLA',
    'RCA',
    'Ticket Status',
    'Comments (Team Follow-up)',
    'Closing Comments',
    'QA Test Results',
    'Tracking Changes'
]

SLIDE_3_KEYS = [
    'Project Name',
    'Review Process',
    'Tool1',
    'Story Point',
    'Acceptance Criteria',
    'If No, please explain_x002e_',
    'Story / Defect Description',
    'Comments (optional)',
    'Module1',
    'Priority1',
    'Owner1',
    'Sprint Tag',
    'Ticket Status1',
    'Comments / Team Follow-up',
    'QA Test Results1',
    'Technical Changes',
    'RCA (Root Cause Analysis)'
]

SLIDE_4_KEYS = [
    'Project Name',
    'PMD Name',
    'US Lead Name',
    'AC Lead Name',
    'Overall Onshore Satisfaction',
    'Overall Client Satisfaction',
    'Are there any active Up-Sell/Cross-Sell opportunities?',
    'Onshore Team Feedback',
    'Client Feedback',
    'Offshore Team Feedback',
    'Up-sell / Cross-sell Opportunities1'
]

SLIDE_5_KEYS = [
    'Project Name',
    'PMD Name',
    'US Lead Name',
    'AC Lead Name',
    'Are DEH tool updates completed and up to date?',
    'Are NexGen Portal updates completed for the project?',
    'If NexGen updates are partial or pending, please provide details_x002e_',
    'Are Peer Reviews conducted regularly?',
    'Are SME Reviews conducted for the project?',
    'If SME reviews are scheduled or not applicable, please provide details_x002e_',
    'Is Concourse being used for RAID logs and documentation?',
    'Are risks, issues, and documents regularly updated in Concourse?',
    'Are Ticket Audits performed regularly?',
    'Is RCA documentation available and up to date (for defects/issues)?',
    'If RCA documentation is missing or pending, please explain_x002e_',
    'Is capacity planning completed and reviewed?',
    'Is vacation tracking and shift roster maintained for the team?'
]

SLIDE_6_KEYS = [
    'Project Name',
    'Operational / Governance',
    'Quality',
    'Environment Health',
    'Risks / Escalations',
    'If Amber or Red is selected for any question, please provide the reason_x002e_'
]

# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================

def get_status_color(status):
    """Return RGB color based on status"""
    if not status:
        return RGBColor(192, 192, 192)  # Gray
    
    status_lower = status.lower()
    
    if 'green' in status_lower or 'yes' in status_lower:
        return RGBColor(146, 208, 80)  # Green
    elif 'amber' in status_lower or 'yellow' in status_lower:
        return RGBColor(255, 192, 0)  # Amber/Orange
    elif 'red' in status_lower or 'no' in status_lower:
        return RGBColor(255, 0, 0)  # Red
    else:
        return RGBColor(217, 217, 217)  # Light Gray for N/A

def add_cell_with_color(table, row, col, text, bg_color=None, font_size=10, bold=False, align_center=False):
    """Add text to a cell with optional background color"""
    cell = table.cell(row, col)
    cell.text = str(text) if text else "N/A"
    
    # Set text formatting
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = 'Calibri'
            run.font.bold = bold
    
    # Set background color if provided
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

def parse_incoming_data(incoming_data):
    """Parse JSON or Python literal string to dict/list"""
    if isinstance(incoming_data, str):
        print("📝 Data received as string, parsing...\n")
        
        try:
            incoming_data = json.loads(incoming_data)
            print("✅ Parsed using JSON\n")
        except json.JSONDecodeError:
            print("⚠️ JSON parsing failed, trying Python literal eval...\n")
            incoming_data = ast.literal_eval(incoming_data)
            print("✅ Parsed using ast.literal_eval\n")
    
    # Convert single dict to list
    if isinstance(incoming_data, dict):
        incoming_data = [incoming_data]
    
    return incoming_data

def calculate_projects_per_slide(total_projects, min_per_slide=3, max_per_slide=5):
    """
    Calculate optimal distribution of projects across slides.
    Ensures each slide has between min_per_slide and max_per_slide projects.
    Returns a list of project counts for each slide.
    
    Args:
        total_projects: Total number of projects to distribute
        min_per_slide: Minimum projects per slide (default 3)
        max_per_slide: Maximum projects per slide (default 5)
    
    Returns:
        List of integers representing projects per slide
    """
    if total_projects == 0:
        return []
    
    if total_projects <= max_per_slide:
        # If total projects fit in one slide, use them all
        return [total_projects]
    
    # Calculate number of slides needed
    num_slides = (total_projects + max_per_slide - 1) // max_per_slide
    
    # Try to distribute evenly
    base_count = total_projects // num_slides
    remainder = total_projects % num_slides
    
    # If base count is less than minimum, reduce number of slides
    if base_count < min_per_slide:
        num_slides = (total_projects + max_per_slide - 1) // max_per_slide
        base_count = total_projects // num_slides
        remainder = total_projects % num_slides
    
    # Create distribution list
    distribution = []
    for i in range(num_slides):
        if i < remainder:
            distribution.append(base_count + 1)
        else:
            distribution.append(base_count)
    
    return distribution

def calculate_column_widths(num_projects, total_width=12.7, fixed_col_width=2.5):
    """
    Calculate dynamic column widths based on number of projects.
    
    Args:
        num_projects: Number of project columns to create
        total_width: Total available width in inches (default 12.7 for standard slide)
        fixed_col_width: Width of the fixed first column (parameter/category names)
    
    Returns:
        Tuple of (fixed_column_width, project_column_width)
    """
    # Calculate remaining width for project columns
    remaining_width = total_width - fixed_col_width
    
    # Divide equally among project columns
    project_col_width = remaining_width / num_projects
    
    return (fixed_col_width, project_col_width)

# ============================================================================
# SLIDE CREATION FUNCTIONS
# ============================================================================

def create_slide_1_project_health(prs, projects_data):
    """
    Slide 1: Project Health Status
    Shows 3-5 projects per slide (dynamically distributed)
    """
    print("📄 Creating Slide 1: Project Health Status...\n")
    
    # Calculate optimal distribution (3-5 projects per slide)
    distribution = calculate_projects_per_slide(len(projects_data), min_per_slide=3, max_per_slide=5)
    
    project_index = 0
    for projects_count in distribution:
        projects_on_slide = projects_data[project_index:project_index + projects_count]
        project_index += projects_count
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title with standardized positioning
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Project Health Status"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.name = 'Arial'
        
        # Create table
        rows = 1 + len(projects_on_slide)  # Header + project rows (fixed: was creating extra row)
        cols = 12
        
        # Standardized margins: 0.5" from edges, ensuring table fits within slide (13.33" × 7.5")
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), 
                                       Inches(12.33), Inches(5.9)).table
        
        # Set column widths
        col_widths = [1.2, 1.1, 0.9, 0.9, 0.8, 0.8, 0.9, 0.9, 1.0, 1.0, 1.5, 1.2]
        for i, col_width in enumerate(col_widths):
            table.columns[i].width = Inches(col_width)
        
        # Header row
        headers = [
            'Project\nName', 'PMD/\nUS Lead/\nAC Lead', 'Contract\nStart Date\n& End Date',
            'Overall\nStatus\n(Red/Amber/\nGreen)', 'Staffing', 'Scope',
            'Project\nGovernance', 'Escalation\nManagement', 'Reason\nfor Amber\n/ Red',
            'Total FTEs', 'Highlights', 'Up-sell/Cross-sell\nOpportunities'
        ]
        
        for col, header in enumerate(headers):
            add_cell_with_color(table, 0, col, header, 
                              bg_color=RGBColor(191, 191, 191), 
                              font_size=9, bold=True, align_center=True)
        
        # Add project data
        for proj_idx, project in enumerate(projects_on_slide):
            row = proj_idx + 1
            
            # Project Name
            add_cell_with_color(table, row, 0, project.get('Project Name', 'N/A'), 
                              font_size=10, bold=True)
            
            # PMD Name
            add_cell_with_color(table, row, 1, project.get('PMD Name', 'N/A'), font_size=9)
            
            # Contract Dates (placeholder)
            add_cell_with_color(table, row, 2, "2021\nSeptember-\n2027 March", font_size=9)
            
            # Status columns with colors
            status_fields = [
                ('Overall Status', 3),
                ('Staffing Status', 4),
                ('Scope Status', 5),
                ('Project Governance Status', 6),
                ('Escalation Management Status1', 7)
            ]
            
            for field_name, col_idx in status_fields:
                status = project.get(field_name, 'N/A')
                status_color = get_status_color(status)
                add_cell_with_color(table, row, col_idx, status, 
                                  bg_color=status_color, font_size=10, 
                                  bold=True, align_center=True)
            
            # Text fields
            add_cell_with_color(table, row, 8, project.get('Reason for Amber / Red', 'N/A'), font_size=9)
            add_cell_with_color(table, row, 9, project.get('Total FTEs', 'N/A'), font_size=9)
            add_cell_with_color(table, row, 10, project.get('Project Highlights', 'N/A'), font_size=9)
            
            # Up-sell opportunities
            upsell_opp = project.get('Up-sell / Cross-sell Opportunities', 'None')
            upsell_details = project.get('Up-sell / Cross-sell Details', '')
            upsell_text = f"{upsell_opp}\n{upsell_details}" if upsell_details and upsell_opp.lower() not in ['none', 'n/a'] else upsell_opp
            add_cell_with_color(table, row, 11, upsell_text, font_size=9)

def create_slide_2_ticket_quality(prs, projects_data):
    """
    Slide 2: Ticket Quality Checks (Incidents)
    Shows 3-5 projects per slide (dynamically distributed)
    """
    print("📄 Creating Slide 2: Ticket Quality Checks (Incidents)...\n")
    
    # Calculate optimal distribution (3-5 projects per slide)
    distribution = calculate_projects_per_slide(len(projects_data), min_per_slide=3, max_per_slide=5)
    
    project_index = 0
    for projects_count in distribution:
        projects_on_slide = projects_data[project_index:project_index + projects_count]
        project_index += projects_count
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title with standardized positioning
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Ticket Quality Checks (Incidents)"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(26)
        title_para.font.bold = True
        title_para.font.name = 'Arial'
        
        # Create table with standardized margins (increased height for 15 rows)
        rows = 14  # 13 parameters + 1 header (fixed: was 15 causing extra blank row)
        cols = 2 + len(projects_on_slide)  # Parameter name + Review Description + Projects
        
        # Create table - dimensions match Delivery Review Summary to fit within slide
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), 
                                       Inches(12.33), Inches(5.9)).table
        
        # Set column widths dynamically - reduced to fit within 12.33" table width
        fixed_col_width, project_col_width = calculate_column_widths(len(projects_on_slide), total_width=12.33, fixed_col_width=5.0)
        table.columns[0].width = Inches(1.8)  # Ticket Quality Parameter (reduced from 2.2)
        table.columns[1].width = Inches(3.2)  # Review Description (reduced from 3.5)
        for i in range(2, cols):
            table.columns[i].width = Inches(project_col_width)  # Dynamic project columns
        
        # Header row with project names
        add_cell_with_color(table, 0, 0, "Ticket Quality Parameter", 
                          bg_color=RGBColor(128, 128, 128), 
                          font_size=11, bold=True, align_center=True)
        add_cell_with_color(table, 0, 1, "Review Description", 
                          bg_color=RGBColor(128, 128, 128), 
                          font_size=11, bold=True, align_center=True)
        
        for proj_idx, project in enumerate(projects_on_slide):
            project_name = project.get('Project Name', f'Project {proj_idx + 1}')
            add_cell_with_color(table, 0, 2 + proj_idx, project_name, 
                              bg_color=RGBColor(128, 128, 128), 
                              font_size=11, bold=True, align_center=True)
        
        # Define parameters and their review questions
        parameters = [
            ('Tool', 'What is the ticketing tool used in project?'),
            ('Description', 'Description clearly added to the ticket?'),
            ('Steps to Reproduce', 'Are the steps to reproduce clearly stated in the ticket?'),
            ('Module', 'Is there a field in the ticketing tool that tags the ticket to particular SE functionality?'),
            ('Priority', 'Is Business Priority captured?'),
            ('Owner', 'Is Owner properly assigned to ticket?'),
            ('SLA', 'Is SLA tracked for the ticket?'),
            ('RCA', 'Is RCA updated on ticket?'),
            ('Ticket Status', 'Is Status being updated regularly and correctly?'),
            ('Comments (Team Follow-up)', 'Is team regularly updating ticket with appropriate and clear comments?'),
            ('Closing Comments', 'Is resolution comments added after issue is fixed?'),
            ('QA Test Results', 'Is QA creating Test task and adding all necessary documentations before closure?'),
            ('Tracking Changes', 'Is the modified components list updated on ticket?')
        ]
        
        # Add parameter rows
        for row_idx, (param_key, review_desc) in enumerate(parameters, start=1):
            # Parameter name (left column)
            add_cell_with_color(table, row_idx, 0, param_key, 
                              bg_color=RGBColor(217, 192, 192), 
                              font_size=10, bold=True)
            
            # Review description (middle column)
            add_cell_with_color(table, row_idx, 1, review_desc, 
                              bg_color=RGBColor(217, 217, 217), 
                              font_size=9)
            
            # Project values
            for proj_idx, project in enumerate(projects_on_slide):
                value = project.get(param_key, 'N/A')
                
                # Determine background color based on value
                if value and str(value).strip().upper() in ['YES', 'GREEN']:
                    bg_color = RGBColor(146, 208, 80)  # Green
                elif value and str(value).strip().upper() in ['NO', 'RED']:
                    bg_color = RGBColor(255, 192, 203)  # Light red
                elif value and str(value).strip().upper() in ['N/A', 'NOT APPLICABLE']:
                    bg_color = RGBColor(217, 217, 217)  # Gray
                else:
                    bg_color = RGBColor(200, 200, 200)  # Default gray
                
                add_cell_with_color(table, row_idx, 2 + proj_idx, value, 
                                  bg_color=bg_color, font_size=9)

def create_slide_3_enhancements(prs, projects_data):
    """
    Slide 3: Quality Checks - Enhancements/Bugs
    Shows 3-5 projects per slide (dynamically distributed)
    """
    print("📄 Creating Slide 3: Quality Checks (Enhancements/Bugs)...\n")
    
    # Calculate optimal distribution (3-5 projects per slide)
    distribution = calculate_projects_per_slide(len(projects_data), min_per_slide=3, max_per_slide=5)
    
    project_index = 0
    for projects_count in distribution:
        projects_on_slide = projects_data[project_index:project_index + projects_count]
        project_index += projects_count
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title with standardized positioning
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Quality Checks (Enhancements/Bugs)"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(26)
        title_para.font.bold = True
        title_para.font.name = 'Arial'
        
        # Create table with standardized margins (increased height for 15 rows)
        rows = 15  # 14 parameters + 1 header
        cols = 2 + len(projects_on_slide)  # Parameter name + Review Description + Projects
        
        # Create table - dimensions match Delivery Review Summary to fit within slide
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), 
                                       Inches(12.33), Inches(5.9)).table
        
        # Set column widths dynamically - reduced to fit within 12.33" table width
        fixed_col_width, project_col_width = calculate_column_widths(len(projects_on_slide), total_width=12.33, fixed_col_width=5.0)
        table.columns[0].width = Inches(1.8)  # Ticket Quality Parameter (reduced from 2.2)
        table.columns[1].width = Inches(3.2)  # Review Description (reduced from 3.5)
        for i in range(2, cols):
            table.columns[i].width = Inches(project_col_width)  # Dynamic project columns
        
        # Header row with project names
        add_cell_with_color(table, 0, 0, "Ticket Quality Parameter", 
                          bg_color=RGBColor(128, 128, 128), 
                          font_size=11, bold=True, align_center=True)
        add_cell_with_color(table, 0, 1, "Review Description", 
                          bg_color=RGBColor(128, 128, 128), 
                          font_size=11, bold=True, align_center=True)
        
        for proj_idx, project in enumerate(projects_on_slide):
            project_name = project.get('Project Name', f'Project {proj_idx + 1}')
            add_cell_with_color(table, 0, 2 + proj_idx, project_name, 
                              bg_color=RGBColor(128, 128, 128), 
                              font_size=11, bold=True, align_center=True)
        
        # Define parameters and their review questions
        parameters = [
            ('Review Process', 'How is the review done?'),
            ('Tool1', 'What is the tool used?'),
            ('Story Point', 'How is the story point measured and assigned?'),
            ('Acceptance Criteria', 'Is AC being captured for stories?'),
            ('Story / Defect Description', 'Is there a proper description added for stories?'),
            ('Module1', 'Is there a field on ticketing tool that tags story to particular functionality?'),
            ('Priority1', 'Is Business Priority captured on story'),
            ('Owner1', 'Is Owner properly assigned to story?'),
            ('Sprint Tag', 'Is correct Sprint tagged to story?'),
            ('Ticket Status1', 'Is Status being updated regularly and correctly?'),
            ('Comments / Team Follow-up', 'Is team regularly following up on ticket if story is blocked?'),
            ('QA Test Results1', 'Is QA creating Test task and adding all necessary documentations before closure?'),
            ('Technical Changes', 'Is the modified components list updated on story?'),
            ('RCA (Root Cause Analysis)', 'Is RCA updated if it is bug?')
        ]
        
        # Add parameter rows
        for row_idx, (param_key, review_desc) in enumerate(parameters, start=1):
            # Parameter name (left column)
            add_cell_with_color(table, row_idx, 0, param_key, 
                              bg_color=RGBColor(217, 192, 192), 
                              font_size=10, bold=True)
            
            # Review description (middle column)
            add_cell_with_color(table, row_idx, 1, review_desc, 
                              bg_color=RGBColor(217, 217, 217), 
                              font_size=9)
            
            # Project values
            for proj_idx, project in enumerate(projects_on_slide):
                value = project.get(param_key, 'N/A')
                
                # Determine background color based on value
                if value and str(value).strip().upper() in ['YES', 'GREEN']:
                    bg_color = RGBColor(146, 208, 80)  # Green
                elif value and str(value).strip().upper() in ['NO', 'RED']:
                    bg_color = RGBColor(255, 192, 203)  # Light red
                elif value and str(value).strip().upper() in ['N/A', 'NOT APPLICABLE']:
                    bg_color = RGBColor(217, 217, 217)  # Gray
                else:
                    bg_color = RGBColor(200, 200, 200)  # Default gray
                
                add_cell_with_color(table, row_idx, 2 + proj_idx, value, 
                                  bg_color=bg_color, font_size=9)

def create_slide_4_feedback_summary(prs, projects_data):
    """
    Slide 4: Feedback Summary
    Shows multiple projects vertically (rows) with feedback columns
    Shows 3-5 projects per slide (dynamically distributed)
    """
    print("📄 Creating Slide 4: Feedback Summary...\n")
    
    # Calculate optimal distribution (3-5 projects per slide)
    distribution = calculate_projects_per_slide(len(projects_data), min_per_slide=3, max_per_slide=5)
    
    project_index = 0
    for projects_count in distribution:
        projects_on_slide = projects_data[project_index:project_index + projects_count]
        project_index += projects_count
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title with standardized positioning
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Feedback Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.name = 'Arial'
        
        # Create table with standardized margins (increased height for multiple rows)
        rows = 1 + len(projects_on_slide)  # Header + project rows
        cols = 8  # Project Name, PMD, US Lead, AC Lead, Onshore Feedback, Client Feedback, Offshore Feedback, Up-sell
        
        # Create table - dimensions match Delivery Review Summary to fit within slide
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), 
                                       Inches(12.33), Inches(5.9)).table
        
        # Set column widths - slightly increased for better readability
        table.columns[0].width = Inches(1.4)   # Project Name (increased from 1.3)
        table.columns[1].width = Inches(1.1)   # PMD (increased from 1.0)
        table.columns[2].width = Inches(1.1)   # US Lead (increased from 1.0)
        table.columns[3].width = Inches(1.1)   # AC Lead (increased from 1.0)
        table.columns[4].width = Inches(2.1)   # Onshore Team Feedback (increased from 2.0)
        table.columns[5].width = Inches(2.1)   # Client Feedback (increased from 2.0)
        table.columns[6].width = Inches(1.9)   # Offshore Team Feedback (increased from 1.8)
        table.columns[7].width = Inches(1.5)   # Up-sell/Cross-sell (increased from 1.4)
        
        # Header row
        headers = [
            'Project Name',
            'PMD',
            'US Lead',
            'AC Lead',
            'Onshore Team\nFeedback',
            'Client Feedback',
            'Offshore Team\nFeedback',
            'Up-Sell/Cross-sell\nOpportunities'
        ]
        
        for col_idx, header in enumerate(headers):
            add_cell_with_color(table, 0, col_idx, header, 
                              bg_color=RGBColor(128, 128, 128), 
                              font_size=10, bold=True, align_center=True)
        
        # Add project data rows
        for proj_idx, project in enumerate(projects_on_slide):
            row = proj_idx + 1
            
            # Project Name
            add_cell_with_color(table, row, 0, project.get('Project Name', 'N/A'), 
                              bg_color=RGBColor(192, 192, 192),
                              font_size=10, bold=True)
            
            # PMD Name
            add_cell_with_color(table, row, 1, project.get('PMD Name', 'N/A'), 
                              bg_color=RGBColor(217, 217, 217),
                              font_size=9)
            
            # US Lead Name
            add_cell_with_color(table, row, 2, project.get('US Lead Name', 'N/A'), 
                              bg_color=RGBColor(217, 217, 217),
                              font_size=9)
            
            # AC Lead Name
            add_cell_with_color(table, row, 3, project.get('AC Lead Name', 'N/A'), 
                              bg_color=RGBColor(217, 217, 217),
                              font_size=9)
            
            # Onshore Team Feedback (with color based on satisfaction)
            onshore_feedback = project.get('Onshore Team Feedback', 'N/A')
            onshore_satisfaction = project.get('Overall Onshore Satisfaction', '')
            onshore_bg = get_feedback_color(onshore_satisfaction, onshore_feedback)
            add_cell_with_color(table, row, 4, onshore_feedback, 
                              bg_color=onshore_bg,
                              font_size=9)
            
            # Client Feedback (with color based on satisfaction)
            client_feedback = project.get('Client Feedback', 'N/A')
            client_satisfaction = project.get('Overall Client Satisfaction', '')
            client_bg = get_feedback_color(client_satisfaction, client_feedback)
            add_cell_with_color(table, row, 5, client_feedback, 
                              bg_color=client_bg,
                              font_size=9)
            
            # Offshore Team Feedback
            offshore_feedback = project.get('Offshore Team Feedback', 'N/A')
            add_cell_with_color(table, row, 6, offshore_feedback, 
                              bg_color=RGBColor(217, 217, 217),
                              font_size=9)
            
            # Up-sell/Cross-sell Opportunities
            upsell_opp = project.get('Up-sell / Cross-sell Opportunities1', 'None')
            upsell_active = project.get('Are there any active Up-Sell/Cross-Sell opportunities?', 'No')
            
            # Color code based on whether there are opportunities
            if upsell_active and str(upsell_active).strip().upper() in ['YES', 'Y']:
                upsell_bg = RGBColor(255, 255, 153)  # Light yellow for opportunities
            else:
                upsell_bg = RGBColor(217, 217, 217)  # Gray for none
            
            add_cell_with_color(table, row, 7, upsell_opp, 
                              bg_color=upsell_bg,
                              font_size=9)

def get_feedback_color(satisfaction_level, feedback_text):
    """
    Determine background color for feedback cells based on satisfaction level or feedback content
    """
    # Check satisfaction level first
    if satisfaction_level:
        sat_lower = str(satisfaction_level).lower()
        if 'high' in sat_lower or 'positive' in sat_lower or 'good' in sat_lower:
            return RGBColor(146, 208, 80)  # Green
        elif 'low' in sat_lower or 'negative' in sat_lower or 'poor' in sat_lower:
            return RGBColor(255, 192, 203)  # Light red
        elif 'medium' in sat_lower or 'neutral' in sat_lower:
            return RGBColor(255, 255, 153)  # Light yellow
    
    # Check feedback text for keywords
    if feedback_text and feedback_text != 'N/A':
        feedback_lower = str(feedback_text).lower()
        
        # Positive keywords
        positive_keywords = ['happy', 'appreciated', 'successful', 'good', 'excellent', 'satisfied']
        if any(keyword in feedback_lower for keyword in positive_keywords):
            return RGBColor(255, 255, 153)  # Light yellow for positive feedback
        
        # Negative keywords
        negative_keywords = ['unhappy', 'issue', 'problem', 'concern', 'disappointed']
        if any(keyword in feedback_lower for keyword in negative_keywords):
            return RGBColor(255, 192, 203)  # Light red for concerns
    
    # Default gray
    return RGBColor(217, 217, 217)

def create_slide_5_salesforce_quality(prs, projects_data):
    """
    Slide 5: Salesforce Project Quality
    Shows multiple projects vertically (rows) with quality metric columns
    Shows 3-5 projects per slide (dynamically distributed)
    """
    print("📄 Creating Slide 5: Salesforce Project Quality...\n")
    
    # Calculate optimal distribution (3-5 projects per slide)
    distribution = calculate_projects_per_slide(len(projects_data), min_per_slide=3, max_per_slide=5)
    
    project_index = 0
    for projects_count in distribution:
        projects_on_slide = projects_data[project_index:project_index + projects_count]
        project_index += projects_count
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title with standardized positioning
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Salesforce Project Quality"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.name = 'Arial'
        
        # Create table with standardized margins (increased height for 12 columns)
        rows = 1 + len(projects_on_slide)  # Header + project rows
        cols = 12  # Project Name, PMD, US Lead, AC Lead, DEH, NexGen, Peer Reviews, SME Reviews, Concourse, Ticket Audits, RCA, Capacity, Vacation
        
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), 
                                       Inches(12.4), Inches(6.2)).table
        
        # Set column widths - increased slightly for better readability
        table.columns[0].width = Inches(1.2)   # Project Name (increased from 1.15)
        table.columns[1].width = Inches(0.95)  # PMD (increased from 0.9)
        table.columns[2].width = Inches(0.95)  # US Lead (increased from 0.9)
        table.columns[3].width = Inches(0.95)  # AC Lead (increased from 0.9)
        table.columns[4].width = Inches(0.95)  # DEH tool Updates (increased from 0.9)
        table.columns[5].width = Inches(1.05)  # NexGen Port Updates (increased from 1.0)
        table.columns[6].width = Inches(0.95)  # Peer Reviews (increased from 0.9)
        table.columns[7].width = Inches(1.05)  # SME Reviews (increased from 1.0)
        table.columns[8].width = Inches(1.1)   # Concourse (increased from 1.05)
        table.columns[9].width = Inches(0.95)  # Ticket Audits (increased from 0.9)
        table.columns[10].width = Inches(1.05) # RCA Documentation (increased from 1.0)
        table.columns[11].width = Inches(1.1)  # Capacity Planning (increased from 1.05)
        
        # Header row
        headers = [
            'Project Name',
            'PMD',
            'US Lead',
            'AC Lead',
            'DEH tool\nUpdates',
            'NexGen\nPort\nUpdates',
            'Peer\nReviews',
            'SME Reviews',
            'Concourse\nAdoption-\nRAID\nLogs/SaveIT',
            'Ticket\nAudits',
            'RCA\nDocument\n-ation',
            'Capacity\nPlanning +\nVacation\ntracker/\nShifts\nRoster'
        ]
        
        for col_idx, header in enumerate(headers):
            add_cell_with_color(table, 0, col_idx, header, 
                              bg_color=RGBColor(128, 128, 128), 
                              font_size=8, bold=True, align_center=True)
        
        # Add project data rows
        for proj_idx, project in enumerate(projects_on_slide):
            row = proj_idx + 1
            
            # Project Name
            add_cell_with_color(table, row, 0, project.get('Project Name', 'N/A'), 
                              bg_color=RGBColor(192, 192, 192),
                              font_size=9, bold=True)
            
            # PMD Name
            add_cell_with_color(table, row, 1, project.get('PMD Name', 'N/A'), 
                              bg_color=RGBColor(217, 217, 217),
                              font_size=8)
            
            # US Lead Name
            add_cell_with_color(table, row, 2, project.get('US Lead Name', 'N/A'), 
                              bg_color=RGBColor(217, 217, 217),
                              font_size=8)
            
            # AC Lead Name
            add_cell_with_color(table, row, 3, project.get('AC Lead Name', 'N/A'), 
                              bg_color=RGBColor(217, 217, 217),
                              font_size=8)
            
            # DEH tool updates
            deh_status = project.get('Are DEH tool updates completed and up to date?', 'N/A')
            deh_bg = get_quality_status_color(deh_status)
            add_cell_with_color(table, row, 4, deh_status, 
                              bg_color=deh_bg,
                              font_size=8, align_center=True)
            
            # NexGen Portal updates
            nexgen_status = project.get('Are NexGen Portal updates completed for the project?', 'N/A')
            nexgen_details = project.get('If NexGen updates are partial or pending, please provide details_x002e_', '')
            
            # If there are details, it means partial/pending
            if nexgen_details and nexgen_details.strip() and nexgen_details.strip().upper() != 'N/A':
                nexgen_display = nexgen_details
                nexgen_bg = RGBColor(255, 192, 0)  # Orange for partial/pending
            else:
                nexgen_display = nexgen_status
                nexgen_bg = get_quality_status_color(nexgen_status)
            
            add_cell_with_color(table, row, 5, nexgen_display, 
                              bg_color=nexgen_bg,
                              font_size=7, align_center=True)
            
            # Peer Reviews
            peer_status = project.get('Are Peer Reviews conducted regularly?', 'N/A')
            peer_bg = get_quality_status_color(peer_status)
            add_cell_with_color(table, row, 6, peer_status, 
                              bg_color=peer_bg,
                              font_size=8, align_center=True)
            
            # SME Reviews
            sme_status = project.get('Are SME Reviews conducted for the project?', 'N/A')
            sme_details = project.get('If SME reviews are scheduled or not applicable, please provide details_x002e_', '')
            
            # If there are details, show them
            if sme_details and sme_details.strip() and sme_details.strip().upper() != 'N/A':
                sme_display = sme_details
                # Check if it's scheduled (orange) or not applicable (gray)
                if 'scheduled' in sme_details.lower() or 'yet' in sme_details.lower():
                    sme_bg = RGBColor(255, 192, 0)  # Orange
                elif 'not applicable' in sme_details.lower() or 'n/a' in sme_details.lower():
                    sme_bg = RGBColor(128, 128, 128)  # Gray
                else:
                    sme_bg = get_quality_status_color(sme_status)
            else:
                sme_display = sme_status
                sme_bg = get_quality_status_color(sme_status)
            
            add_cell_with_color(table, row, 7, sme_display, 
                              bg_color=sme_bg,
                              font_size=7, align_center=True)
            
            # Concourse (combined check)
            concourse_used = project.get('Is Concourse being used for RAID logs and documentation?', 'N/A')
            concourse_updated = project.get('Are risks, issues, and documents regularly updated in Concourse?', 'N/A')
            
            # Combine both checks
            if str(concourse_used).strip().upper() == 'YES' and str(concourse_updated).strip().upper() == 'YES':
                concourse_display = 'Concourse available\nRisks, documents\nare updated'
                concourse_bg = RGBColor(0, 176, 80)  # Green
            elif str(concourse_used).strip().upper() == 'YES':
                concourse_display = 'Concourse available\nNeeds to be\nupdated'
                concourse_bg = RGBColor(255, 192, 0)  # Orange
            else:
                concourse_display = concourse_used
                concourse_bg = get_quality_status_color(concourse_used)
            
            add_cell_with_color(table, row, 8, concourse_display, 
                              bg_color=concourse_bg,
                              font_size=7, align_center=True)
            
            # Ticket Audits
            ticket_status = project.get('Are Ticket Audits performed regularly?', 'N/A')
            ticket_bg = get_quality_status_color(ticket_status)
            add_cell_with_color(table, row, 9, ticket_status, 
                              bg_color=ticket_bg,
                              font_size=8, align_center=True)
            
            # RCA Documentation
            rca_status = project.get('Is RCA documentation available and up to date (for defects/issues)?', 'N/A')
            rca_details = project.get('If RCA documentation is missing or pending, please explain_x002e_', '')
            
            # If there are details, show them
            if rca_details and rca_details.strip() and rca_details.strip().upper() != 'N/A':
                rca_display = rca_details
                rca_bg = RGBColor(255, 192, 0)  # Orange for missing/pending
            else:
                rca_display = rca_status
                rca_bg = get_quality_status_color(rca_status)
            
            add_cell_with_color(table, row, 10, rca_display, 
                              bg_color=rca_bg,
                              font_size=7, align_center=True)
            
            # Capacity Planning and Vacation tracking (combined)
            capacity_status = project.get('Is capacity planning completed and reviewed?', 'N/A')
            vacation_status = project.get('Is vacation tracking and shift roster maintained for the team?', 'N/A')
            
            # Both should be Yes for green
            if str(capacity_status).strip().upper() == 'YES' and str(vacation_status).strip().upper() == 'YES':
                combined_display = 'Yes'
                combined_bg = RGBColor(0, 176, 80)  # Green
            elif str(capacity_status).strip().upper() == 'YES' or str(vacation_status).strip().upper() == 'YES':
                combined_display = 'Partial'
                combined_bg = RGBColor(255, 192, 0)  # Orange
            else:
                combined_display = 'N/A'
                combined_bg = RGBColor(128, 128, 128)  # Gray
            
            add_cell_with_color(table, row, 11, combined_display, 
                              bg_color=combined_bg,
                              font_size=8, align_center=True)

def get_quality_status_color(status):
    """
    Determine background color for quality status cells
    """
    if not status:
        return RGBColor(128, 128, 128)  # Gray for empty
    
    status_str = str(status).strip().upper()
    
    # Green for Yes/Complete/Done
    if status_str in ['YES', 'Y', 'COMPLETE', 'COMPLETED', 'DONE', 'UP TO DATE']:
        return RGBColor(0, 176, 80)  # Green
    
    # Orange/Yellow for Partial/Pending/In Progress
    elif status_str in ['PARTIAL', 'PENDING', 'IN PROGRESS', 'SCHEDULED', 'ONGOING']:
        return RGBColor(255, 192, 0)  # Orange
    
    # Gray for N/A/Not Applicable/No
    elif status_str in ['N/A', 'NA', 'NOT APPLICABLE', 'NO', 'NONE']:
        return RGBColor(128, 128, 128)  # Gray
    
    # Check for keywords in longer text
    status_lower = status.lower()
    if 'not applicable' in status_lower or 'n/a' in status_lower:
        return RGBColor(128, 128, 128)  # Gray
    elif 'pending' in status_lower or 'partial' in status_lower or 'missing' in status_lower:
        return RGBColor(255, 192, 0)  # Orange
    elif 'yes' in status_lower or 'complete' in status_lower or 'available' in status_lower:
        return RGBColor(0, 176, 80)  # Green
    
    # Default gray for unknown
    return RGBColor(128, 128, 128)

def create_slide_6_delivery_review(prs, projects_data):
    """
    Slide 6: Delivery Review Summary
    Shows projects as columns (horizontal) with category rows
    Displays status indicators (Green/Amber dots) for each category
    Shows 3-5 projects per slide (dynamically distributed)
    """
    print("📄 Creating Slide 6: Delivery Review Summary...\n")
    
    # Calculate optimal distribution (3-5 projects per slide)
    distribution = calculate_projects_per_slide(len(projects_data), min_per_slide=3, max_per_slide=5)
    
    project_index = 0
    for projects_count in distribution:
        projects_on_slide = projects_data[project_index:project_index + projects_count]
        project_index += projects_count
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title with standardized positioning
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Delivery Review Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.name = 'Arial'
        
        # Create table for status indicators with standardized margins
        # Rows: Header + 4 categories (Operational, Quality, Environment, Risks) + Reasons row
        rows = 6
        cols = 1 + len(projects_on_slide)  # Category column + Project columns
        
        # Create table - dimensions match standard to fit within slide boundaries
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), 
                                       Inches(12.33), Inches(5.9)).table
        
        # Set column widths dynamically based on number of projects
        fixed_col_width, project_col_width = calculate_column_widths(len(projects_on_slide), total_width=12.33, fixed_col_width=2.5)
        table.columns[0].width = Inches(fixed_col_width)  # Category/Project Name column
        for i in range(1, cols):
            table.columns[i].width = Inches(project_col_width)  # Dynamic project columns
        
        # Header row - Project Names
        add_cell_with_color(table, 0, 0, "Project Name", 
                          bg_color=RGBColor(128, 128, 128), 
                          font_size=11, bold=True, align_center=True)
        
        for proj_idx, project in enumerate(projects_on_slide):
            project_name = project.get('Project Name', f'Project {proj_idx + 1}')
            add_cell_with_color(table, 0, 1 + proj_idx, project_name, 
                              bg_color=RGBColor(128, 128, 128), 
                              font_size=11, bold=True, align_center=True)
        
        # Category rows with sub-items
        categories = [
            ('Operational/Governance', [
                'Staffing',
                'Scope',
                'Project Governance',
                'Escalation Management'
            ], 'Operational / Governance'),
            ('Quality', [
                'Peer review tracker',
                'SME Reviews',
                'Ticket Audits',
                'RCA Documentation for Incidents',
                'Capacity Planning',
                'Vacation tracker/ Shifts Roaster',
                'DEH tool updates',
                'NexGen Port updates'
            ], 'Quality'),
            ('Environment Health', [
                'Structured Environment Pipeline',
                'Sandbox refresh frequency',
                'Environments maintenance',
                'Version control'
            ], 'Environment Health'),
            ('Risks / Escalations', [
                'Issues ( P1)',
                'Challenges',
                'Risks',
                'Mitigation',
                'RAID logs'
            ], 'Risks / Escalations')
        ]
        
        # Add category rows (rows 1-4)
        for row_idx, (category_display, sub_items, data_key) in enumerate(categories, start=1):
            # Category name with sub-items
            category_text = f"{category_display}\n" + "\n".join([f"  • {item}" for item in sub_items])
            
            add_cell_with_color(table, row_idx, 0, category_text, 
                              bg_color=RGBColor(188, 143, 143), 
                              font_size=8, bold=True)
            
            # Add status indicators for each project
            for proj_idx, project in enumerate(projects_on_slide):
                status = project.get(data_key, 'N/A')
                
                # Determine status color and display
                status_display = get_status_indicator(status)
                status_bg = RGBColor(217, 217, 217)  # Light gray background
                
                cell = table.cell(row_idx, 1 + proj_idx)
                cell.text = status_display
                cell.fill.solid()
                cell.fill.fore_color.rgb = status_bg
                
                # Center align and format
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(24)  # Large font for the dot indicator
                        run.font.name = 'Calibri'
                        
                        # Color the text (dot) based on status
                        if 'green' in status.lower() or status.strip().upper() == 'YES':
                            run.font.color.rgb = RGBColor(0, 176, 80)  # Green
                        elif 'amber' in status.lower() or 'yellow' in status.lower():
                            run.font.color.rgb = RGBColor(255, 192, 0)  # Amber/Orange
                        elif 'red' in status.lower():
                            run.font.color.rgb = RGBColor(255, 0, 0)  # Red
                        else:
                            run.font.color.rgb = RGBColor(128, 128, 128)  # Gray
                
                cell.text_frame.vertical_anchor = 1  # Middle alignment
        
        # Reasons for Amber/Red status row (row 5)
        add_cell_with_color(table, 5, 0, "Reasons for Amber status", 
                          bg_color=RGBColor(188, 143, 143), 
                          font_size=10, bold=True)
        
        for proj_idx, project in enumerate(projects_on_slide):
            reason = project.get('If Amber or Red is selected for any question, please provide the reason_x002e_', 'NA')
            
            # Format reason as bullet points if it's a longer text
            if reason and reason.strip() and reason.strip().upper() != 'NA':
                reason_display = reason
                reason_bg = RGBColor(222, 184, 135)  # Tan/beige for reasons
            else:
                reason_display = 'NA'
                reason_bg = RGBColor(217, 217, 217)  # Gray
            
            add_cell_with_color(table, 5, 1 + proj_idx, reason_display, 
                              bg_color=reason_bg,
                              font_size=8)

def get_status_indicator(status):
    """
    Return a visual indicator (dot) based on status
    """
    if not status:
        return '●'  # Bullet point (will be colored gray)
    
    status_lower = status.lower()
    
    if 'green' in status_lower or status.strip().upper() == 'YES':
        return '●'  # Green dot
    elif 'amber' in status_lower or 'yellow' in status_lower:
        return '●'  # Amber dot
    elif 'red' in status_lower:
        return '●'  # Red dot
    else:
        return '●'  # Default dot

def create_cover_slide(prs):
    """
    Create professional cover slide with modern design
    Based on the provided template with dark gray and orange color scheme
    """
    print("📄 Creating Cover Slide...\n")
    
    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Get current date
    current_date = datetime.now().strftime('%B %Y')
    
    # Add dark gray background rectangle (left side - covers about 70% width)
    dark_gray_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(9.3), Inches(7.5)
    )
    dark_gray_bg.fill.solid()
    dark_gray_bg.fill.fore_color.rgb = RGBColor(68, 68, 68)  # Dark gray
    dark_gray_bg.line.fill.background()
    
    # Add orange accent rectangle (right side - covers about 30% width)
    orange_accent = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(9.3), Inches(0),
        Inches(4.03), Inches(5.1)
    )
    orange_accent.fill.solid()
    orange_accent.fill.fore_color.rgb = RGBColor(214, 124, 34)  # Orange
    orange_accent.line.fill.background()
    
    # Add orange-red date bar at bottom
    date_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(5.1),
        Inches(9.3), Inches(1.3)
    )
    date_bar.fill.solid()
    date_bar.fill.fore_color.rgb = RGBColor(184, 71, 42)  # Orange-red
    date_bar.line.fill.background()
    
    # Add title text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.0), Inches(2.0))
    title_frame = title_box.text_frame
    title_frame.text = "XYZ Managed Services Portfolio\nDelivery Overview"
    title_frame.word_wrap = True
    
    # Format title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.name = 'Georgia'  # Changed from Arial to Georgia
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        paragraph.alignment = PP_ALIGN.LEFT
    
    # Add date text
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(6.0), Inches(0.6))
    date_frame = date_box.text_frame
    date_frame.text = current_date
    
    # Format date
    for paragraph in date_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.bold = False
        paragraph.font.name = 'Georgia'  # Changed from Arial to Georgia
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        paragraph.alignment = PP_ALIGN.LEFT

def create_section_divider(prs, section_title):
    """
    Create section divider slide with light gray background and orange-red accent bar
    Based on the provided template design
    
    Args:
        prs: Presentation object
        section_title: Title of the section to display
    """
    print(f"📄 Creating Section Divider: {section_title}...\n")
    
    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add white background for entire slide (changed from light gray)
    white_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    white_bg.fill.solid()
    white_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    white_bg.line.fill.background()
    
    # Add orange-red accent bar at bottom
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(6.6),
        Inches(13.33), Inches(0.9)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(184, 71, 42)  # Orange-red
    accent_bar.line.fill.background()
    
    # Add section title text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.0), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_frame.word_wrap = True
    
    # Format title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(48)
        paragraph.font.bold = True
        paragraph.font.name = 'Georgia'  # Changed from Arial to Georgia
        paragraph.font.color.rgb = RGBColor(68, 68, 68)  # Dark gray text
        paragraph.alignment = PP_ALIGN.LEFT

# ============================================================================
# MAIN PPT GENERATION FUNCTION
# ============================================================================

def create_complete_presentation(projects_data):
    """
    Create complete PowerPoint presentation with all slides
    """
    print("="*80)
    print("🎨 GENERATING COMPLETE POWERPOINT PRESENTATION")
    print("="*80 + "\n")
    
    # Create new presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    
    
    # Generate all slide types in the requested order
    # 0. Cover Slide
    create_cover_slide(prs)
    
    # 1. Delivery Review Summary Section
    create_section_divider(prs, "Delivery Review Summary")
    create_slide_6_delivery_review(prs, projects_data)
    
    # 2. Project Health Status Section
    create_section_divider(prs, "Project Health Status")
    create_slide_1_project_health(prs, projects_data)
    
    # 3. Salesforce Project Quality Section
    create_section_divider(prs, "Salesforce Project Quality")
    create_slide_5_salesforce_quality(prs, projects_data)
    
    # 4. Ticket Quality Checks - Incidents Section
    create_section_divider(prs, "Ticket Quality Checks - Incidents")
    create_slide_2_ticket_quality(prs, projects_data)
    
    # 5. Quality Checks - Enhancements/Bugs Section
    create_section_divider(prs, "Quality Checks - Enhancements/Bugs")
    create_slide_3_enhancements(prs, projects_data)
    
    # 6. Feedback Summary Section
    create_section_divider(prs, "Feedback Summary")
    create_slide_4_feedback_summary(prs, projects_data)
    
    # Add more slide creation functions here as you develop them
    # create_slide_7_xxx(prs, projects_data)
    # etc.
    
    
    # Save presentation to memory instead of disk
    output_filename = f"Complete_Project_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    
    # Create in-memory bytes buffer
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    # Encode to base64 for Power Automate
    ppt_base64 = base64.b64encode(ppt_buffer.read()).decode('utf-8')
    
    print("="*80)
    print(f"✅ POWERPOINT CREATED SUCCESSFULLY: {output_filename}")
    print("="*80 + "\n")
    
    return output_filename, ppt_base64

# ============================================================================
# SERVERLESS FUNCTION ENTRY POINT
# ============================================================================

def main(args):
    """
    Main entry point for DigitalOcean Serverless Function
    
    Args:
        args: dict or list containing project data
        
    Returns:
        dict with status, message, and results
    """
    print("\n" + "="*80)
    print("SUCCESS! Data received")
    print("="*80 + "\n")
    
    try:
        # Parse the incoming data
        incoming_data = parse_incoming_data(args)
        
        if not isinstance(incoming_data, list):
            raise ValueError("Data must be a list of dictionaries")
        
        print(f"📊 Total number of project records: {len(incoming_data)}\n")
        
        # Print extracted data for verification
        print("="*80)
        print("EXTRACTED DATA PREVIEW")
        print("="*80 + "\n")
        
        for index, project in enumerate(incoming_data[:2], start=1):  # Show first 2 projects
            print(f"PROJECT #{index}: {project.get('Project Name', 'Unknown')}")
            print(f"  Status: {project.get('Overall Status', 'N/A')}")
            print(f"  Tool: {project.get('Tool', 'N/A')}")
            print()
        
        # Create complete PowerPoint presentation
        ppt_filename, ppt_base64 = create_complete_presentation(incoming_data)
        
        # Send PowerPoint to Power Automate
        print("📤 Sending PowerPoint to Power Automate...\n")
        
        payload = {
            "filename": ppt_filename,
            "fileContent": ppt_base64
        }
        
        try:
            response = requests.post(
                POWER_AUTOMATE_ENDPOINT,
                json=payload,
                headers={'Content-Type': 'application/json'},
                timeout=30
            )
            response.raise_for_status()
            
            print("="*80)
            print("✅ POWERPOINT SENT TO POWER AUTOMATE SUCCESSFULLY!")
            print("="*80 + "\n")
            
            # Success response
            return {
                "status": "Success", 
                "message": "PowerPoint created and sent to Power Automate successfully!",
                "records_processed": len(incoming_data),
                "ppt_filename": ppt_filename,
                "power_automate_status": "Sent"
            }
            
        except requests.exceptions.RequestException as e:
            print(f"❌ ERROR: Failed to send to Power Automate")
            print(f"Error: {e}")
            return {
                "status": "Partial Success",
                "message": "PowerPoint created but failed to send to Power Automate",
                "records_processed": len(incoming_data),
                "ppt_filename": ppt_filename,
                "error": str(e)
            }
    
    except (json.JSONDecodeError, ValueError, SyntaxError) as e:
        print(f"❌ ERROR: Failed to parse data")
        print(f"Error: {e}")
        traceback.print_exc()
        return {
            "status": "Error", 
            "message": f"Parsing failed: {str(e)}"
        }
    
    except Exception as e:
        print(f"❌ ERROR: Unexpected error occurred")
        print(f"Error: {e}")
        traceback.print_exc()
        return {
            "status": "Error", 
            "message": f"Processing failed: {str(e)}"
        }
